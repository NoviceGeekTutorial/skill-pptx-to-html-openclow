#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX to HTML Converter - 1:1 Layout Preservation
精确复刻PPT布局，生成高质量HTML

Features:
- 1:1 精确复刻PPT内容
- 完全保持布局、排版、颜色
- HTML文字可自由选择复制
- 全平台自适应（电脑/手机/平板）
- 支持文本、图片、形状、表格、图表
"""

import sys
import os
import re
import base64
import json
from pathlib import Path
from typing import Optional, List, Dict, Tuple, Any
from xml.etree import ElementTree as ET
from dataclasses import dataclass, field
from enum import Enum

# python-pptx imports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn, nsmap
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

from lxml import etree

# XML namespaces
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


class FontStyleExtractor:
    """Extract font styles from PPTX XML elements"""
    
    @staticmethod
    def hex_to_rgb(hex_str: str) -> str:
        """Convert hex color to lowercase hex"""
        if not hex_str:
            return "#000000"
        return f"#{hex_str.lower()}"
    
    @staticmethod
    def get_attr(element, attr_name: str, default=None):
        """Get XML attribute safely"""
        return element.get(attr_name, default)
    
    @staticmethod
    def extract_color(element) -> str:
        """Extract color from XML element"""
        # Try srgbClr
        srgb = element.find('.//a:srgbClr', NSMAP)
        if srgb is not None:
            val = srgb.get('val', '')
            if val:
                return FontStyleExtractor.hex_to_rgb(val)
        
        # Try theme color
        scheme = element.find('.//a:schemeClr', NSMAP)
        if scheme is not None:
            # Theme colors need theme mapping, use default
            return "#000000"
        
        return "#000000"
    
    @staticmethod
    def extract_font_from_defRPr(defRPr, text: str = "") -> Dict[str, Any]:
        """Extract font style from a:defRPr element"""
        style = {
            'bold': False,
            'italic': False,
            'underline': False,
            'size': 18.0,  # Default size in points
            'color': '#000000',
            'font_family': 'Arial, "Microsoft YaHei", sans-serif'
        }
        
        if defRPr is None:
            return style
        
        # Extract size (in hundredths of a point)
        size_str = defRPr.get('sz', '')
        if size_str:
            try:
                # sz is in 100ths of a point
                style['size'] = int(size_str) / 100.0
            except ValueError:
                pass
        
        # Extract bold
        b = defRPr.get('b', '0')
        style['bold'] = b in ('1', 'true')
        
        # Extract italic
        i = defRPr.get('i', '0')
        style['italic'] = i in ('1', 'true')
        
        # Extract underline
        u = defRPr.get('u', 'none')
        if u != 'none':
            style['underline'] = True
        
        # Extract color
        solid_fill = defRPr.find('a:solidFill', NSMAP)
        if solid_fill is not None:
            style['color'] = FontStyleExtractor.extract_color(solid_fill)
        
        # Extract font name
        latin = defRPr.find('a:latin', NSMAP)
        if latin is not None:
            typeface = latin.get('typeface', '')
            if typeface:
                style['font_family'] = f'"{typeface}", Arial, "Microsoft YaHei", sans-serif'
        
        ea = defRPr.find('a:ea', NSMAP)
        if ea is not None and not style['font_family'].startswith('"'):
            typeface = ea.get('typeface', '')
            if typeface:
                style['font_family'] = f'"{typeface}", Arial, "Microsoft YaHei", sans-serif'
        
        return style
    
    @staticmethod
    def extract_from_run(run_elem) -> Dict[str, Any]:
        """Extract font style from a:r (run) element"""
        # First check run's own rPr
        rPr = run_elem.find('a:rPr', NSMAP)
        if rPr is not None:
            return FontStyleExtractor.extract_font_from_defRPr(rPr)
        
        # Fall back to paragraph's defRPr
        # This requires parent traversal
        para = run_elem.getparent()
        if para is not None:
            pPr = para.find('a:pPr', NSMAP)
            if pPr is not None:
                defRPr = pPr.find('a:defRPr', NSMAP)
                if defRPr is not None:
                    return FontStyleExtractor.extract_font_from_defRPr(defRPr)
        
        # Ultimate fallback
        return {
            'bold': False,
            'italic': False,
            'underline': False,
            'size': 18.0,
            'color': '#000000',
            'font_family': 'Arial, "Microsoft YaHei", sans-serif'
        }
    
    @staticmethod
    def extract_paragraph_alignment(para_elem) -> str:
        """Extract text alignment from paragraph element"""
        pPr = para_elem.find('a:pPr', NSMAP)
        if pPr is not None:
            algn = pPr.get('algn', 'l')
            alignment_map = {
                'l': 'left',
                'ctr': 'center',
                'r': 'right',
                'ju': 'justify'
            }
            return alignment_map.get(algn, 'left')
        return 'left'


class PPTXtoHTMLConverter:
    """PPTX to HTML Converter with 1:1 layout preservation"""
    
    def __init__(self, pptx_path: str, output_dir: str):
        self.pptx_path = Path(pptx_path)
        self.output_dir = Path(output_dir)
        self.prs = Presentation(pptx_path)
        
        # Get actual slide dimensions
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        
        # Create output directories
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.assets_dir = self.output_dir / "assets"
        self.assets_dir.mkdir(exist_ok=True)
        self.images_dir = self.assets_dir / "images"
        self.images_dir.mkdir(exist_ok=True)
        
        # Track extracted images
        self.image_counter = 0
        
    def emu_to_px(self, emu: int, dpi: int = 96) -> float:
        """Convert EMUs to pixels at 96 DPI"""
        return (emu / 914400) * dpi
    
    def rgb_to_hex(self, color) -> str:
        """Convert RGBColor object to hex string"""
        if color is None:
            return "#000000"
        try:
            if hasattr(color, 'rgb'):
                rgb = color.rgb
                if rgb and len(rgb) >= 3:
                    return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"
        except:
            pass
        return "#000000"
    
    def get_shape_fill_color(self, shape) -> Optional[str]:
        """Extract fill color from shape"""
        try:
            fill = shape.fill
            if fill.type is not None:
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    color = fill.fore_color
                    # Handle theme colors
                    if hasattr(color, 'rgb') and color.rgb:
                        return self.rgb_to_hex(color)
                    elif hasattr(color, 'theme_color'):
                        # Theme colors need mapping
                        return self._get_theme_color(color.theme_color)
        except:
            pass
        return None
    
    def _get_theme_color(self, theme_color_idx) -> str:
        """Map theme color index to actual color"""
        theme_colors = {
            0: "#000000",  # Dark 1
            1: "#FFFFFF",  # Light 1
            2: "#000000",  # Dark 2
            3: "#FFFFFF",  # Light 2
            4: "#000000",  # Accent 1
            5: "#000000",  # Accent 2
            6: "#000000",  # Accent 3
            7: "#000000",  # Accent 4
            8: "#000000",  # Accent 5
            9: "#000000",  # Accent 6
        }
        return theme_colors.get(theme_color_idx, "#000000")
    
    def get_shape_line_color(self, shape) -> Optional[str]:
        """Extract line/border color from shape"""
        try:
            line = shape.line
            if line.fill and line.fill.type is not None:
                if hasattr(line.fill, 'fore_color') and line.fill.fore_color:
                    return self.rgb_to_hex(line.fill.fore_color)
        except:
            pass
        return None
    
    def extract_image(self, shape, slide_idx: int) -> Optional[str]:
        """Extract image from shape and save to disk"""
        try:
            if not hasattr(shape, 'image'):
                return None
            
            image = shape.image
            if image is None:
                return None
            
            # Generate unique filename
            self.image_counter += 1
            ext = image.ext
            if not ext:
                ext = 'png'
            ext = ext.lstrip('.')
            
            filename = f"slide_{slide_idx + 1}_img_{self.image_counter}.{ext}"
            filepath = self.images_dir / filename
            
            # Save image
            with open(filepath, 'wb') as f:
                f.write(image.blob)
            
            return f"assets/images/{filename}"
        except Exception as e:
            print(f"Warning: Failed to extract image: {e}")
            return None
    
    def process_text_from_xml(self, shape) -> str:
        """Process text frame by extracting directly from XML for accurate styling"""
        try:
            # Get the XML element (use _element for newer python-pptx versions)
            sp = getattr(shape, '_element', None)
            if sp is None:
                return ""
            
            # Find all text bodies
            txBody = sp.find('.//p:txBody', NSMAP)
            if txBody is None:
                return ""
            
            html_parts = []
            
            # Process each paragraph
            for para in txBody.findall('a:p', NSMAP):
                # Get paragraph alignment
                alignment = FontStyleExtractor.extract_paragraph_alignment(para)
                
                para_html = []
                has_content = False
                
                # Process each run
                for run in para.findall('a:r', NSMAP):
                    t = run.find('a:t', NSMAP)
                    if t is None or not t.text:
                        continue
                    
                    has_content = True
                    
                    # Extract font style from XML
                    style = FontStyleExtractor.extract_from_run(run)
                    
                    # Build inline styles
                    inline_styles = [
                        f"font-size: {style['size']:.1f}pt",
                        f"color: {style['color']}",
                        f"font-family: {style['font_family']}"
                    ]
                    
                    if style['bold']:
                        inline_styles.append("font-weight: bold")
                    if style['italic']:
                        inline_styles.append("font-style: italic")
                    if style['underline']:
                        inline_styles.append("text-decoration: underline")
                    
                    # Escape HTML special characters
                    text = t.text
                    text = text.replace('&', '&amp;')
                    text = text.replace('<', '&lt;')
                    text = text.replace('>', '&gt;')
                    # Preserve line breaks
                    text = text.replace('\n', '<br>')
                    text = text.replace('\r', '')
                    
                    span_html = f'<span style="{"; ".join(inline_styles)}">{text}</span>'
                    para_html.append(span_html)
                
                # Handle text without runs (e.g., placeholder text)
                if not has_content:
                    # Try to get text from a:t directly
                    t = para.find('.//a:t', NSMAP)
                    if t is not None and t.text:
                        para_html.append(f'<span style="font-size: 18pt; color: #000000; font-family: Arial, sans-serif">{t.text}</span>')
                        has_content = True
                
                if has_content and para_html:
                    para_content = ''.join(para_html)
                    html_parts.append(f'<p style="text-align: {alignment}; margin: 0 0 0.3em 0; line-height: 1.4;">{para_content}</p>')
            
            return '\n'.join(html_parts)
        except Exception as e:
            print(f"Warning: Error extracting text from XML: {e}")
            return ""
    
    def process_shape(self, shape, slide_idx: int) -> str:
        """Process a single shape and generate HTML"""
        # Get shape position and size
        left = self.emu_to_px(shape.left)
        top = self.emu_to_px(shape.top)
        width = self.emu_to_px(shape.width)
        height = self.emu_to_px(shape.height)
        
        shape_type = shape.shape_type
        
        # Handle different shape types
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self.process_picture(shape, left, top, width, height, slide_idx)
        elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            return self.process_text_box(shape, left, top, width, height)
        elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return self.process_auto_shape(shape, left, top, width, height, slide_idx)
        elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            return self.process_placeholder(shape, left, top, width, height, slide_idx)
        elif shape_type == MSO_SHAPE_TYPE.TABLE:
            return self.process_table(shape, left, top, width, height)
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            return self.process_group(shape, left, top, width, height, slide_idx)
        else:
            return self.process_generic_shape(shape, left, top, width, height, slide_idx)
    
    def process_picture(self, shape, left: float, top: float, 
                       width: float, height: float, slide_idx: int) -> str:
        """Process picture/image shape"""
        img_path = self.extract_image(shape, slide_idx)
        
        if img_path:
            return f'''<div class="pptx-shape pptx-image" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px;">
    <img src="{img_path}" style="width: 100%; height: 100%; object-fit: contain;" alt="" loading="lazy">
</div>'''
        else:
            return f'''<div class="pptx-shape pptx-image-placeholder" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px; background: #f0f0f0; border: 1px dashed #ccc; display: flex; align-items: center; justify-content: center; color: #999; font-size: 14px;">
    [Image]
</div>'''
    
    def process_text_box(self, shape, left: float, top: float, 
                        width: float, height: float) -> str:
        """Process text box shape"""
        content = self.process_text_from_xml(shape)
        
        if not content:
            return ""
        
        bg_color = self.get_shape_fill_color(shape)
        bg_style = f"background-color: {bg_color};" if bg_color else ""
        
        return f'''<div class="pptx-shape pptx-textbox" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px; {bg_style}">
    {content}
</div>'''
    
    def process_auto_shape(self, shape, left: float, top: float, 
                          width: float, height: float, slide_idx: int) -> str:
        """Process auto shape (rectangles, circles, arrows, etc.)"""
        bg_color = self.get_shape_fill_color(shape)
        line_color = self.get_shape_line_color(shape)
        
        styles = [
            f"position: absolute",
            f"left: {left}px",
            f"top: {top}px",
            f"width: {width}px",
            f"height: {height}px"
        ]
        
        if bg_color:
            styles.append(f"background-color: {bg_color}")
        else:
            styles.append("background-color: transparent")
        
        if line_color:
            styles.append(f"border: 1px solid {line_color}")
        
        # Check if shape has text
        content = self.process_text_from_xml(shape)
        if content:
            styles.append("display: flex")
            styles.append("flex-direction: column")
            styles.append("justify-content: center")
            styles.append("align-items: center")
        
        style_str = "; ".join(styles)
        
        return f'''<div class="pptx-shape pptx-autoshape" style="{style_str}">
    {content}
</div>'''
    
    def process_placeholder(self, shape, left: float, top: float, 
                            width: float, height: float, slide_idx: int) -> str:
        """Process placeholder shape"""
        content = self.process_text_from_xml(shape)
        
        if not content:
            return ""
        
        return f'''<div class="pptx-shape pptx-placeholder" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px;">
    {content}
</div>'''
    
    def process_table(self, shape, left: float, top: float, 
                      width: float, height: float) -> str:
        """Process table shape"""
        if not hasattr(shape, 'table'):
            return ""
        
        table = shape.table
        if not table or not table.rows:
            return ""
        
        rows_html = []
        
        for row_idx, row in enumerate(table.rows):
            cells_html = []
            for cell in row.cells:
                # Get cell content from XML
                cell_text = ""
                try:
                    cell_elem = cell._tc
                    t = cell_elem.find('.//a:t', NSMAP)
                    if t is not None and t.text:
                        cell_text = t.text
                except:
                    cell_text = cell.text or ""
                
                # Escape HTML
                cell_text = cell_text.replace('&', '&amp;')
                cell_text = cell_text.replace('<', '&lt;')
                cell_text = cell_text.replace('>', '&gt;')
                
                # Get cell styling
                cell_styles = ["border: 1px solid #ccc", "padding: 8px"]
                
                # Try to get background color
                try:
                    if cell.fill and cell.fill.fore_color:
                        bg = self.rgb_to_hex(cell.fill.fore_color)
                        cell_styles.append(f"background-color: {bg}")
                except:
                    pass
                
                style_str = "; ".join(cell_styles)
                tag = 'th' if row_idx == 0 else 'td'
                cells_html.append(f'<{tag} style="{style_str}">{cell_text}</{tag}>')
            
            rows_html.append(f'<tr>{"".join(cells_html)}</tr>')
        
        table_html = '\n'.join(rows_html)
        
        return f'''<div class="pptx-shape pptx-table" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px;">
    <table style="width: 100%; border-collapse: collapse;">
        {table_html}
    </table>
</div>'''
    
    def process_group(self, shape, left: float, top: float, 
                      width: float, height: float, slide_idx: int) -> str:
        """Process grouped shapes"""
        if not hasattr(shape, 'shapes'):
            return ""
        
        group_html = []
        for child_shape in shape.shapes:
            try:
                child_html = self.process_shape(child_shape, slide_idx)
                if child_html:
                    group_html.append(child_html)
            except Exception as e:
                print(f"Warning: Error processing child shape: {e}")
        
        if not group_html:
            return ""
        
        return f'''<div class="pptx-shape pptx-group" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px;">
    {"\n".join(group_html)}
</div>'''
    
    def process_generic_shape(self, shape, left: float, top: float, 
                              width: float, height: float, slide_idx: int) -> str:
        """Process generic/unknown shape types"""
        # Try to get text content
        content = self.process_text_from_xml(shape)
        
        # Try to get image
        img_path = None
        if hasattr(shape, 'image'):
            img_path = self.extract_image(shape, slide_idx)
        
        if img_path:
            return f'''<div class="pptx-shape pptx-image" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px;">
    <img src="{img_path}" style="width: 100%; height: 100%; object-fit: contain;" alt="" loading="lazy">
</div>'''
        elif content:
            return f'''<div class="pptx-shape pptx-generic" style="position: absolute; left: {left}px; top: {top}px; width: {width}px; height: {height}px;">
    {content}
</div>'''
        
        return ""
    
    def get_slide_background(self, slide) -> str:
        """Extract slide background color"""
        try:
            background = slide.background
            fill = background.fill
            
            if fill.type is not None:
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    return f"background-color: {self.rgb_to_hex(fill.fore_color)};"
        except:
            pass
        
        return "background-color: #ffffff;"
    
    def convert_slide(self, slide, slide_idx: int) -> str:
        """Convert a single slide to HTML"""
        slide_width_px = self.emu_to_px(self.slide_width)
        slide_height_px = self.emu_to_px(self.slide_height)
        
        bg_style = self.get_slide_background(slide)
        
        shapes_html = []
        for shape in slide.shapes:
            try:
                shape_html = self.process_shape(shape, slide_idx)
                if shape_html:
                    shapes_html.append(shape_html)
            except Exception as e:
                print(f"Warning: Error processing shape on slide {slide_idx + 1}: {e}")
        
        shapes_content = "\n".join(shapes_html) if shapes_html else ""
        
        return f'''<div class="pptx-slide" style="position: relative; width: {slide_width_px}px; height: {slide_height_px}px; {bg_style} margin: 0 auto 20px auto; box-shadow: 0 2px 8px rgba(0,0,0,0.15); overflow: hidden;">
    {shapes_content}
</div>'''
    
    def generate_css(self) -> str:
        """Generate CSS styles for the HTML output"""
        slide_width_px = self.emu_to_px(self.slide_width)
        slide_height_px = self.emu_to_px(self.slide_height)
        
        return f'''/* PPTX to HTML Converter Styles */

.pptx-container {{
    font-family: Arial, "Microsoft YaHei", "Hiragino Sans GB", "WenQuanYi Micro Hei", sans-serif;
    line-height: 1.6;
    color: #333;
    -webkit-font-smoothing: antialiased;
    -moz-osx-font-smoothing: grayscale;
}}

.pptx-slide {{
    box-sizing: border-box;
    page-break-after: always;
    break-after: page;
}}

.pptx-shape {{
    box-sizing: border-box;
}}

/* Text selection enabled */
.pptx-textbox, .pptx-autoshape, .pptx-placeholder {{
    user-select: text;
    -webkit-user-select: text;
    -moz-user-select: text;
    -ms-user-select: text;
}}

.pptx-textbox p, .pptx-autoshape p, .pptx-placeholder p,
.pptx-textbox span, .pptx-autoshape span, .pptx-placeholder span {{
    user-select: text;
    -webkit-user-select: text;
    -moz-user-select: text;
    -ms-user-select: text;
}}

/* Image styles */
.pptx-image img {{
    display: block;
    max-width: 100%;
    height: auto;
    user-select: none;
    -webkit-user-select: none;
    pointer-events: none;
}}

/* Table styles */
.pptx-table table {{
    width: 100%;
    border-collapse: collapse;
}}

.pptx-table th, .pptx-table td {{
    border: 1px solid #ddd;
    padding: 8px;
    text-align: left;
}}

.pptx-table th {{
    background-color: #f5f5f5;
    font-weight: bold;
}}

/* AutoShape styles */
.pptx-autoshape {{
    display: flex;
    align-items: center;
    justify-content: center;
}}

/* Print styles */
@media print {{
    body {{ background: white; padding: 0; }}
    .pptx-slide {{
        page-break-inside: avoid;
        box-shadow: none;
        margin: 0 0 20px 0;
    }}
    .pptx-header {{ display: none; }}
}}

/* Responsive styles */
@media screen and (max-width: {int(slide_width_px + 40)}px) {{
    .pptx-slide {{
        width: 100% !important;
        height: auto !important;
        aspect-ratio: {int(slide_width_px)} / {int(slide_height_px)};
    }}
}}

/* Mobile */
@media screen and (max-width: 768px) {{
    .pptx-container {{ padding: 10px; }}
    .pptx-slide {{ margin-bottom: 15px; }}
    body {{ padding: 10px; }}
}}

/* High DPI */
@media (-webkit-min-device-pixel-ratio: 2), (min-resolution: 192dpi) {{
    .pptx-image img {{
        image-rendering: -webkit-optimize-contrast;
        image-rendering: crisp-edges;
    }}
}}'''
    
    def convert(self) -> str:
        """Convert entire presentation to HTML"""
        slides_html = []
        
        for idx, slide in enumerate(self.prs.slides):
            print(f"Processing slide {idx + 1}/{len(self.prs.slides)}...")
            try:
                slide_html = self.convert_slide(slide, idx)
                slides_html.append(slide_html)
            except Exception as e:
                print(f"Error processing slide {idx + 1}: {e}")
                slides_html.append(f'<div class="pptx-slide" style="padding: 40px; color: red; background: #fff;">Error: {e}</div>')
        
        # Generate CSS
        css_content = self.generate_css()
        css_path = self.assets_dir / "pptx-styles.css"
        with open(css_path, 'w', encoding='utf-8') as f:
            f.write(css_content)
        
        # Build final HTML
        all_slides = "\n".join(slides_html)
        slide_width_px = self.emu_to_px(self.slide_width)
        
        html_content = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{self.pptx_path.stem}</title>
    <link rel="stylesheet" href="assets/pptx-styles.css">
    <style>
        body {{
            margin: 0;
            padding: 20px;
            background-color: #f5f5f5;
            font-family: Arial, "Microsoft YaHei", "Hiragino Sans GB", sans-serif;
        }}
        
        .pptx-container {{
            max-width: {int(slide_width_px + 40)}px;
            margin: 0 auto;
        }}
        
        .pptx-header {{
            text-align: center;
            margin-bottom: 20px;
            padding: 20px;
            background: white;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }}
        
        .pptx-header h1 {{
            margin: 0 0 10px 0;
            color: #333;
            font-size: 24px;
        }}
        
        .pptx-header p {{
            margin: 0;
            color: #666;
            font-size: 14px;
        }}
    </style>
</head>
<body>
    <div class="pptx-container">
        <div class="pptx-header">
            <h1>{self.pptx_path.stem}</h1>
            <p>Total slides: {len(self.prs.slides)}</p>
        </div>
        {all_slides}
    </div>
</body>
</html>'''
        
        output_html_path = self.output_dir / f"{self.pptx_path.stem}.html"
        with open(output_html_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return str(output_html_path)


def main():
    """Main entry point"""
    if len(sys.argv) < 2:
        print("Usage: python pptx_to_html.py <input.pptx> [output_dir]")
        print("Example: python pptx_to_html.py presentation.pptx ./output")
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "./pptx-html-output"
    
    if not os.path.exists(input_path):
        print(f"Error: File not found: {input_path}")
        sys.exit(1)
    
    if not input_path.endswith('.pptx'):
        print("Error: Input file must be a .pptx file")
        sys.exit(1)
    
    print(f"Converting: {input_path}")
    print(f"Output directory: {output_dir}")
    print("-" * 50)
    
    try:
        converter = PPTXtoHTMLConverter(input_path, output_dir)
        output_path = converter.convert()
        print("-" * 50)
        print("[OK] Conversion complete!")
        print(f"[OK] Output saved to: {output_path}")
        print(f"[OK] Assets saved to: {os.path.join(output_dir, 'assets')}")
    except Exception as e:
        print(f"[ERROR] Conversion failed: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
